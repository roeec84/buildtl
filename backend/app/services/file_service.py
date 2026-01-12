"""
File upload and processing service.
Handles document uploads and text extraction for RAG.
"""
import os
import aiofiles
from pathlib import Path
from typing import List, Optional
from fastapi import UploadFile
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
from bs4 import BeautifulSoup
from app.core.config import settings


class FileService:
    """Service for handling file uploads and text extraction."""

    SUPPORTED_EXTENSIONS = {
        '.pdf', '.txt', '.md',
        '.docx', '.doc',
        '.pptx', '.ppt',
        '.xlsx', '.xls', '.csv',
        '.html', '.htm'
    }

    def __init__(self):
        """Initialize file service and create upload directory."""
        self.upload_dir = Path(settings.UPLOAD_DIR)
        self.upload_dir.mkdir(parents=True, exist_ok=True)

    async def save_upload_file(
        self,
        upload_file: UploadFile,
        username: str
    ) -> str:
        """
        Save an uploaded file to disk.

        Args:
            upload_file: The uploaded file
            username: Username for file organization

        Returns:
            File path where the file was saved

        Raises:
            ValueError: If file type is not supported
        """
        # Validate file extension
        file_ext = Path(upload_file.filename).suffix.lower()
        if file_ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type: {file_ext}. "
                f"Supported types: {', '.join(self.SUPPORTED_EXTENSIONS)}"
            )

        # Create user directory
        user_dir = self.upload_dir / username
        user_dir.mkdir(exist_ok=True)

        # Save file
        file_path = user_dir / upload_file.filename

        async with aiofiles.open(file_path, 'wb') as f:
            content = await upload_file.read()
            await f.write(content)

        return str(file_path)

    async def extract_text(self, file_path: str) -> str:
        """
        Extract text from a file.

        Args:
            file_path: Path to the file

        Returns:
            Extracted text content

        Raises:
            ValueError: If file type is not supported
        """
        file_ext = Path(file_path).suffix.lower()

        if file_ext == '.pdf':
            return await self._extract_from_pdf(file_path)
        elif file_ext == '.txt' or file_ext == '.md':
            return await self._extract_from_text(file_path)
        elif file_ext in ['.docx', '.doc']:
            return await self._extract_from_docx(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return await self._extract_from_pptx(file_path)
        elif file_ext in ['.xlsx', '.xls', '.csv']:
            return await self._extract_from_excel(file_path)
        elif file_ext in ['.html', '.htm']:
            return await self._extract_from_html(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")

    async def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file."""
        text = []
        with open(file_path, 'rb') as file:
            pdf_reader = pypdf.PdfReader(file)
            for page in pdf_reader.pages:
                text.append(page.extract_text())
        return '\n'.join(text)

    async def _extract_from_text(self, file_path: str) -> str:
        """Extract text from plain text file."""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
            return await file.read()

    async def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from Word document."""
        doc = DocxDocument(file_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return '\n'.join(text)

    async def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation."""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

    async def _extract_from_excel(self, file_path: str) -> str:
        """Extract text from Excel file."""
        if file_path.endswith('.csv'):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)

        # Convert DataFrame to text
        return df.to_string()

    async def _extract_from_html(self, file_path: str) -> str:
        """Extract text from HTML file."""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
            html_content = await file.read()

        soup = BeautifulSoup(html_content, 'lxml')
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Get text
        text = soup.get_text()

        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)

        return text

    async def chunk_text(
        self,
        text: str,
        chunk_size: int = 1000,
        chunk_overlap: int = 200
    ) -> List[str]:
        """
        Split text into chunks for vector storage.

        Args:
            text: Text to split
            chunk_size: Size of each chunk
            chunk_overlap: Overlap between chunks

        Returns:
            List of text chunks
        """
        chunks = []
        start = 0

        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            start = end - chunk_overlap

        return chunks
